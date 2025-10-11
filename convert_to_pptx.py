import markdown
from pptx import Presentation
from pptx.util import Inches

def md_to_pptx(md_file, pptx_file):
    """
    Converts a Markdown file to a PowerPoint presentation.
    """
    with open(md_file, 'r') as f:
        md_content = f.read()

    prs = Presentation()
    slides = md_content.split('\n---\n')

    for slide_md in slides:
        if slide_md.strip():
            parts = slide_md.split('\n\n')
            title = parts[0].replace('#', '').strip()
            content = '\n\n'.join(parts[1:]).strip()

            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            
            title_shape = slide.shapes.title
            title_shape.text = title
            
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.text = content

    prs.save(pptx_file)

if __name__ == '__main__':
    md_to_pptx('/Users/Bhavesh/Documents/GitHub/ports_digital_twin/planning/presentation_21Sep/presentation_slides_walkthrough_v5(11 Oct 2025).md', 'presentation.pptx')